"""
Document processor for handling various file types and extracting text.
"""

import os
import re
import logging
import importlib
from typing import Dict, List, Optional, Tuple, Any, Union, BinaryIO
from abc import ABC, abstractmethod

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentProcessor(ABC):
    """Base abstract class for document processors."""
    
    @abstractmethod
    def extract_text(self, file_path: str) -> str:
        """
        Extract text content from a document.
        
        Args:
            file_path: Path to the document file.
            
        Returns:
            Extracted text content.
        """
        pass
    
    @abstractmethod
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from a document.
        
        Args:
            file_path: Path to the document file.
            
        Returns:
            Dictionary of metadata.
        """
        pass
    
    @staticmethod
    def get_processor_for_file(file_path: str) -> 'DocumentProcessor':
        """
        Factory method to get the appropriate processor for a file.
        
        Args:
            file_path: Path to the document file.
            
        Returns:
            An instance of a DocumentProcessor subclass.
            
        Raises:
            ValueError: If the file type is not supported.
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == '.pdf':
            return PDFProcessor()
        elif ext in ('.docx', '.doc'):
            return DocxProcessor()
        elif ext in ('.xlsx', '.xls'):
            return ExcelProcessor()
        elif ext == '.pptx':
            return PowerPointProcessor()
        elif ext in ('.txt', '.md', '.py', '.java', '.js', '.html', '.css', '.json', '.xml'):
            return TextProcessor()
        elif ext in ('.htm', '.html'):
            return HTMLProcessor()
        elif ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif'):
            return ImageProcessor()
        else:
            return TextProcessor()  # Fallback to text processor


class TextProcessor(DocumentProcessor):
    """Processor for plain text files."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with a different encoding if UTF-8 fails
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get basic file metadata."""
        return {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }


class PDFProcessor(DocumentProcessor):
    """Processor for PDF files."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from PDF files."""
        try:
            # Try to import PyPDF2
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                
                # Extract text from each page
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n\n"
                
                return text
        except ImportError:
            logger.warning("PyPDF2 is not installed. Cannot process PDF files.")
            return f"[PDF TEXT EXTRACTION FAILED: PyPDF2 not installed. Please install with 'pip install PyPDF2']"
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return f"[PDF TEXT EXTRACTION ERROR: {str(e)}]"
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get PDF metadata."""
        metadata = {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }
        
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                pdf_info = pdf_reader.metadata
                
                if pdf_info:
                    for key, value in pdf_info.items():
                        # Clean up the key name (remove leading '/')
                        clean_key = key[1:] if key.startswith('/') else key
                        metadata[clean_key] = value
                
                metadata['pages'] = len(pdf_reader.pages)
                
        except ImportError:
            metadata['error'] = "PyPDF2 not installed"
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata


class DocxProcessor(DocumentProcessor):
    """Processor for Microsoft Word (DOCX) files."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            import docx
            
            doc = docx.Document(file_path)
            text = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                text.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            text.append(para.text)
            
            return '\n'.join(text)
        except ImportError:
            logger.warning("python-docx is not installed. Cannot process DOCX files.")
            return f"[DOCX TEXT EXTRACTION FAILED: python-docx not installed. Please install with 'pip install python-docx']"
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return f"[DOCX TEXT EXTRACTION ERROR: {str(e)}]"
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get DOCX metadata."""
        metadata = {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }
        
        try:
            import docx
            
            doc = docx.Document(file_path)
            core_props = doc.core_properties
            
            metadata['author'] = core_props.author
            metadata['created'] = core_props.created
            metadata['last_modified_by'] = core_props.last_modified_by
            metadata['modified'] = core_props.modified
            metadata['title'] = core_props.title
            metadata['subject'] = core_props.subject
            metadata['keywords'] = core_props.keywords
            metadata['comments'] = core_props.comments
            metadata['category'] = core_props.category
            
        except ImportError:
            metadata['error'] = "python-docx not installed"
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata


class ExcelProcessor(DocumentProcessor):
    """Processor for Microsoft Excel files."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from Excel files."""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text = []
            
            # Extract text from each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                
                # Extract cell values
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text:
                        text.append('\t'.join(row_text))
            
            return '\n'.join(text)
        except ImportError:
            logger.warning("openpyxl is not installed. Cannot process Excel files.")
            return f"[EXCEL TEXT EXTRACTION FAILED: openpyxl not installed. Please install with 'pip install openpyxl']"
        except Exception as e:
            logger.error(f"Error extracting text from Excel: {str(e)}")
            return f"[EXCEL TEXT EXTRACTION ERROR: {str(e)}]"
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get Excel metadata."""
        metadata = {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }
        
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            metadata['sheets'] = workbook.sheetnames
            metadata['sheet_count'] = len(workbook.sheetnames)
            
            # Try to get document properties
            if hasattr(workbook, 'properties'):
                props = workbook.properties
                metadata['title'] = props.title
                metadata['creator'] = props.creator
                metadata['subject'] = props.subject
                metadata['description'] = props.description
                metadata['keywords'] = props.keywords
                metadata['category'] = props.category
                metadata['last_modified_by'] = props.lastModifiedBy
            
        except ImportError:
            metadata['error'] = "openpyxl not installed"
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata


class PowerPointProcessor(DocumentProcessor):
    """Processor for Microsoft PowerPoint files."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = []
            
            # Extract text from each slide
            for i, slide in enumerate(prs.slides):
                text.append(f"Slide {i+1}:")
                
                # Extract text from shapes (textboxes, titles, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            
            return '\n'.join(text)
        except ImportError:
            logger.warning("python-pptx is not installed. Cannot process PowerPoint files.")
            return f"[POWERPOINT TEXT EXTRACTION FAILED: python-pptx not installed. Please install with 'pip install python-pptx']"
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint: {str(e)}")
            return f"[POWERPOINT TEXT EXTRACTION ERROR: {str(e)}]"
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get PowerPoint metadata."""
        metadata = {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            metadata['slide_count'] = len(prs.slides)
            
            # Try to get document properties
            core_props = prs.core_properties
            if core_props:
                metadata['author'] = core_props.author
                metadata['title'] = core_props.title
                metadata['subject'] = core_props.subject
                metadata['keywords'] = core_props.keywords
                metadata['comments'] = core_props.comments
                metadata['category'] = core_props.category
            
        except ImportError:
            metadata['error'] = "python-pptx not installed"
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata


class HTMLProcessor(DocumentProcessor):
    """Processor for HTML files."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from HTML files."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'lxml')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            # Get text
            text = soup.get_text(separator=' ', strip=True)
            
            # Break into lines and remove leading and trailing space on each
            lines = (line.strip() for line in text.splitlines())
            
            # Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            
            # Join the non-empty lines
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        except ImportError:
            logger.warning("BeautifulSoup and/or lxml are not installed. Cannot process HTML files.")
            return f"[HTML TEXT EXTRACTION FAILED: BeautifulSoup/lxml not installed. Please install with 'pip install beautifulsoup4 lxml']"
        except Exception as e:
            logger.error(f"Error extracting text from HTML: {str(e)}")
            return f"[HTML TEXT EXTRACTION ERROR: {str(e)}]"
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get HTML metadata."""
        metadata = {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }
        
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'lxml')
            
            # Get title
            if soup.title:
                metadata['title'] = soup.title.string
            
            # Get meta tags
            for meta in soup.find_all('meta'):
                if meta.get('name'):
                    metadata[meta.get('name')] = meta.get('content')
            
        except ImportError:
            metadata['error'] = "BeautifulSoup/lxml not installed"
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata


class ImageProcessor(DocumentProcessor):
    """Processor for image files with OCR capabilities."""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from images using OCR."""
        try:
            import pytesseract
            from PIL import Image
            
            # Open image
            image = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            return text
        except ImportError:
            logger.warning("pytesseract and/or Pillow are not installed. Cannot process images with OCR.")
            return f"[IMAGE OCR FAILED: pytesseract/Pillow not installed. Please install with 'pip install pytesseract Pillow']"
        except Exception as e:
            logger.error(f"Error extracting text from image: {str(e)}")
            return f"[IMAGE OCR ERROR: {str(e)}]"
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get image metadata."""
        metadata = {
            'filename': os.path.basename(file_path),
            'path': os.path.abspath(file_path),
            'size': os.path.getsize(file_path),
            'modified': os.path.getmtime(file_path),
            'created': os.path.getctime(file_path)
        }
        
        try:
            from PIL import Image
            
            # Open image
            image = Image.open(file_path)
            
            # Get image details
            metadata['format'] = image.format
            metadata['mode'] = image.mode
            metadata['width'] = image.width
            metadata['height'] = image.height
            
            # Get EXIF data if available
            if hasattr(image, '_getexif') and callable(image._getexif):
                exif = image._getexif()
                if exif:
                    for tag_id, value in exif.items():
                        # Convert tag ID to name if possible
                        from PIL.ExifTags import TAGS
                        tag = TAGS.get(tag_id, tag_id)
                        metadata[f'exif_{tag}'] = value
            
        except ImportError:
            metadata['error'] = "Pillow not installed"
        except Exception as e:
            metadata['error'] = str(e)
        
        return metadata